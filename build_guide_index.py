# build_guide_index.py
#
# ONE-TIME setup script. Run this whenever the user manuals change.
# It reads the two PowerPoint manuals, breaks them into one chunk per
# slide, asks OpenAI to turn each chunk into an "embedding" (a list of
# numbers that captures the slide's MEANING), and saves everything to
# guide_index.json.
#
# That JSON file is the whole "index" the chatbot's guide tool will later
# search. It contains ONLY user-guide text -- no client data, no asset
# rows, nothing private -- so it is safe to send to OpenAI and safe to
# keep as a plain file.
#
# Why embeddings instead of keyword search: users ask in everyday words
# ("how do I count my inventory") but the manual uses jargon ("Stock
# Taking"). Embeddings match on meaning, so the right slide is found even
# when the words are completely different. That word-gap is the entire
# reason this chatbot exists, so it's the reason we embed.

import json
import os
import re

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

load_dotenv()
client = OpenAI()  # same OpenAI setup / .env key as the rest of the project

# text-embedding-3-small: cheap, small, and more than good enough here.
# Embedding all ~254 slides costs a fraction of a cent, one time.
EMBEDDING_MODEL = "text-embedding-3-small"

# The two manuals, each tagged with a "source" label so a mobile question
# can later be answered from mobile slides only, and vice versa.
MANUALS = [
    {"path": "T-Cube_Web_User_Manual__For_Clients_.pptx", "source": "web"},
    {"path": "Appendix_F_-_Tcube_Mobile_User_Manual.pptx", "source": "mobile"},
]

OUTPUT_FILE = "guide_index.json"

# OPTIONAL extra knowledge file. The manuals explain each screen on its
# own, but some important things are never written down on any single
# slide -- for example that a Stock Take record must be created in the WEB
# system before the mobile app can scan against it. RAG can only retrieve
# what exists, so those cross-module / cross-platform facts have to be
# written down somewhere to be findable. Put them in guide_notes.md and
# they get indexed right alongside the slides.
#
# Format: separate each note with a line containing only "---".
# (If the file doesn't exist, it's simply skipped.)
NOTES_FILE = "guide_notes.md"

# Slides whose text is basically front-matter / navigation, not real
# instructions. We skip these so a real question never retrieves the
# table of contents or the copyright notice.
SKIP_MARKERS = (
    "table of contents",
    "revision history",
    "copyright notice",
    "copyright and disclaimer",
)

# Phrases that mark the COVER / contact slide -- company address, postal
# details, "Updated By", the "Page of XXX" placeholder. A slide is only
# dropped for these if it's ALSO short (see _is_noise), so a real content
# slide that happens to mention the company name is safe. We deliberately
# do NOT drop "Overview" slides: those actually explain what TCube is and
# answer real "what is this system?" questions.
COVER_MARKERS = (
    "boon lay way",        # the company street address
    "tradehub",            # building name on the cover
    "updated by:",         # cover metadata
    "date updated:",       # cover metadata
)


def _slide_text(slide) -> str:
    """Pull all the readable text off one slide (titles, body, table cells)
    and join it into a single clean string. Images are ignored -- we only
    index words, not pictures."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if line:
                    parts.append(line)
        if shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
    text = "\n".join(parts)
    # Collapse boilerplate that repeats on every slide.
    text = re.sub(r"\bUSER REFERENCE\b", "", text)
    text = re.sub(r"\bCONFIDENTIAL\b", "", text)
    text = re.sub(r"Page\s*\d+\s*of\s*\w+", "", text)
    return text.strip()


def _is_noise(text: str) -> bool:
    """True if a slide is front-matter (TOC, revision list, copyright,
    cover/contact page) or just too short to be a useful instruction."""
    low = text.lower()
    if any(marker in low for marker in SKIP_MARKERS):
        return True
    if len(text.split()) < 8:
        return True
    # Cover / contact slide: has address-or-metadata markers AND is short.
    # The length gate is what protects real content -- an instructional
    # slide is long and won't be dropped even if it names the company.
    if len(text.split()) < 45 and any(marker in low for marker in COVER_MARKERS):
        return True
    return False


def build_note_chunks() -> list[dict]:
    """Load hand-written notes from guide_notes.md, if it exists.

    These fill gaps the manuals leave implicit. They are tagged with
    source 'note' so you can tell in the search results that an answer
    came from your own notes rather than from a manual slide."""
    if not os.path.exists(NOTES_FILE):
        print(f"  notes  : {NOTES_FILE} not found (skipping -- this is fine)")
        return []

    with open(NOTES_FILE, "r", encoding="utf-8") as f:
        raw = f.read()

    # Strip <!-- ... --> comments: the file's own header instructions are
    # written as a comment so they don't get indexed as if they were
    # knowledge about TCube.
    raw = re.sub(r"<!--.*?-->", "", raw, flags=re.DOTALL)

    chunks = []
    # Split on a line containing only "---" so each note is its own chunk.
    for i, block in enumerate(re.split(r"^\s*---\s*$", raw, flags=re.MULTILINE), start=1):
        text = block.strip()
        if len(text.split()) < 5:      # skip blank / trivial blocks
            continue
        chunks.append({
            "source": "note",
            "slide": i,                 # note number, not a real slide
            "text": text,
        })
    print(f"  notes  : loaded {len(chunks)} note chunks from {NOTES_FILE}")
    return chunks


def build_chunks() -> list[dict]:
    """Turn both manuals into a flat list of chunks, one per real slide."""
    chunks = []
    for manual in MANUALS:
        prs = Presentation(manual["path"])
        kept = 0
        skipped = 0
        for slide_number, slide in enumerate(prs.slides, start=1):
            text = _slide_text(slide)
            if _is_noise(text):
                skipped += 1
                continue
            chunks.append({
                "source": manual["source"],
                "slide": slide_number,
                "text": text,
            })
            kept += 1
        print(f"  {manual['source']:6} : kept {kept} slides, skipped {skipped} noise slides")

    # Add your own notes on top of the manual slides.
    chunks.extend(build_note_chunks())
    return chunks


def embed_chunks(chunks: list[dict]) -> list[dict]:
    """Ask OpenAI for an embedding for every chunk's text, in batches."""
    BATCH = 100
    for start in range(0, len(chunks), BATCH):
        batch = chunks[start:start + BATCH]
        response = client.embeddings.create(
            model=EMBEDDING_MODEL,
            input=[c["text"] for c in batch],
        )
        for chunk, item in zip(batch, response.data):
            chunk["embedding"] = item.embedding
        print(f"  embedded chunks {start + 1}-{start + len(batch)}")
    return chunks


def main():
    print("Step 1: reading manuals and building chunks...")
    chunks = build_chunks()
    print(f"  total chunks: {len(chunks)}\n")

    print("Step 2: embedding chunks via OpenAI...")
    chunks = embed_chunks(chunks)
    print()

    print(f"Step 3: saving to {OUTPUT_FILE}...")
    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump(chunks, f)
    print(f"  done -- {len(chunks)} chunks saved.\n")

    print("=" * 60)
    print("REVIEW SAMPLE (first 3 of each source):")
    print("=" * 60)
    # Counts start empty and get created on demand, so this keeps working
    # no matter what sources exist (web, mobile, note, or any added later).
    shown: dict[str, int] = {}
    for c in chunks:
        source = c["source"]
        if shown.get(source, 0) < 3:
            shown[source] = shown.get(source, 0) + 1
            preview = c["text"].replace("\n", " ")
            if len(preview) > 200:
                preview = preview[:200] + "..."
            print(f"\n[{source} {c['slide']}] {preview}")
    print()


if __name__ == "__main__":
    main()